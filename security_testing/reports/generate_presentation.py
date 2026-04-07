#!/usr/bin/env python3
"""Generate a 15-minute PowerPoint presentation for the smoltcp security testing project."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

FIGURES = os.path.join(os.path.dirname(__file__), 'figures')
OUT = os.path.join(os.path.dirname(__file__), 'smoltcp_presentation.pptx')

# ── Colors ──
DARK = RGBColor(0x1E, 0x29, 0x3B)
BLUE = RGBColor(0x1E, 0x40, 0xAF)
LIGHT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
GRAY = RGBColor(0x64, 0x74, 0x8B)
ACCENT = RGBColor(0x8B, 0x5C, 0xF6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=DARK):
    """Fill slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height, font_size=18, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_slide(slide, title, bullets, img_path=None):
    """Standard content slide with title, bullets, and optional image."""
    add_bg(slide, DARK)
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    add_text(slide, title, 0.5, 0.15, 12, 0.9, font_size=32, color=WHITE, bold=True)

    # Bullets
    img_width = 5.5 if img_path else 0
    text_width = 11.5 if not img_path else 6.5

    for i, bullet in enumerate(bullets):
        add_text(slide, f"• {bullet}", 0.7, 1.6 + i * 0.65, text_width, 0.6,
                font_size=18, color=LIGHT)

    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(7.5), Inches(1.5), Inches(5.5))


def add_image_slide(slide, title, img_path, caption=""):
    """Full-width image slide."""
    add_bg(slide, DARK)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    add_text(slide, title, 0.5, 0.15, 12, 0.9, font_size=32, color=WHITE, bold=True)

    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.5), Inches(10.3), Inches(5.3))

    if caption:
        add_text(slide, caption, 1.5, 6.9, 10, 0.5, font_size=14, color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, DARK)

# Gradient-like top bar
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_BLUE
shape.line.fill.background()

add_text(slide, "Quality Assessment & Security Testing", 1, 1.5, 11, 1.2,
         font_size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "smoltcp v0.12.0 — Embedded TCP/IP Stack", 1, 2.8, 11, 0.8,
         font_size=28, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, "Comprehensive Security Analysis & Quality Assurance Report", 1, 3.8, 11, 0.6,
         font_size=20, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, "April 2026", 1, 5, 11, 0.6,
         font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

# Bottom accent
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.35), prs.slide_width, Inches(0.15))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_BLUE
shape.line.fill.background()


# ══════════════════════════════════════════════════════════════════
# SLIDE 2: Agenda
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "Agenda", [
    "Project Overview & Scope",
    "Mutation Adequacy Testing (333 Mutants)",
    "Input Space Partitioning (696 Test Cases)",
    "Security Fuzzing with libFuzzer",
    "RFC Conformance Testing",
    "Performance Benchmarking",
    "White-Box Coverage Analysis",
    "Software Safety Assessment",
    "Cross-Platform Portability",
    "Key Findings & Recommendations",
])


# ══════════════════════════════════════════════════════════════════
# SLIDE 3: Project Overview
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "Project Overview", [
    "smoltcp v0.12.0 — standalone, event-driven TCP/IP stack in Rust",
    "Targets embedded, bare-metal, and resource-constrained systems",
    "Zero-copy packet parsing with Rust's memory safety guarantees",
    "Focus modules: assembler.rs, ipv4.rs, udp.rs, tcp.rs",
    "Tested on Windows 10 (MSVC) and WSL Ubuntu 24.04 (GNU/LLVM)",
    "10 quality dimensions assessed across dual platforms",
    "Custom automation scripts for reproducible, cross-platform execution",
])


# ══════════════════════════════════════════════════════════════════
# SLIDE 4: Mutation Overview
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "Mutation Adequacy Testing", [
    "333 first-order mutants across 4 source files",
    "7 mutation operator classes: const, arith, rel, bool, bit, assign, shift",
    "Custom Python engine with module-scoped test filtering",
    "Per-mutant execution time: ~3.4 seconds",
    "Results identical on Windows and WSL (deterministic)",
], os.path.join(FIGURES, 'fig2_mutation_pie.png'))


# ══════════════════════════════════════════════════════════════════
# SLIDE 5: Mutation Results
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_image_slide(slide, "Mutation Kill Ratio by Module",
                os.path.join(FIGURES, 'fig1_mutation_kill_ratio.png'),
                "assembler.rs: 89.41% | ipv4.rs: 61.36% | udp.rs: 57.50% | tcp.rs: 63.33%")


# ══════════════════════════════════════════════════════════════════
# SLIDE 6: Equivalent Mutant Analysis
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "Equivalent Mutant Analysis", [
    "5 surviving mutants analyzed in detail",
    "2 classified as equivalent (debug assertions, display-only branches)",
    "3 classified as killable (boundary conditions in add/remove_front)",
    "Adjusted adequacy: K/(M−E) = 229/331 = 69.18%",
    "Targeted tests recommended for overlapping segment assembly",
], os.path.join(FIGURES, 'fig13_equivalent_mutants.png'))


# ══════════════════════════════════════════════════════════════════
# SLIDE 7: Operator Distribution
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_image_slide(slide, "Mutation Operator Distribution",
                os.path.join(FIGURES, 'fig3_operator_distribution.png'),
                "Constant replacements dominate (98 mutants), followed by arithmetic (52) and relational (48)")


# ══════════════════════════════════════════════════════════════════
# SLIDE 8: Input Space Partitioning
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "Input Space Partitioning", [
    "UDP parsing: 6 dimensions × full combinatorial = 480 test cases",
    "Dimensions: buf_len, len_field, dst_port, checksum, ip_family, rx_on",
    "IPv4 combinatorial: 216 additional test cases",
    "Total: 696 test cases — 100% pass rate on both platforms",
    "Key finding: checksum=0 accepted for UDP over IPv6",
], os.path.join(FIGURES, 'fig7_partition_heatmap.png'))


# ══════════════════════════════════════════════════════════════════
# SLIDE 9: Security Fuzzing
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "Security Fuzzing with libFuzzer", [
    "5 fuzz targets: packet_parser, tcp_headers, dhcp, 802.15.4, 6LoWPAN",
    "60s initial + 180s extended campaigns per target",
    "759 unique coverage edges discovered (packet_parser)",
    "Zero crash artifacts across all targets",
    "Rust's memory safety prevents buffer overflows and null derefs",
    "Windows blocked by MSVC/ASAN DLL incompatibility",
], os.path.join(FIGURES, 'fig9_fuzzing_coverage.png'))


# ══════════════════════════════════════════════════════════════════
# SLIDE 10: Conformance Testing
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_image_slide(slide, "RFC Conformance Testing",
                os.path.join(FIGURES, 'fig6_conformance.png'),
                "9/9 cases matched expected behavior | 2 specification deviations documented")


# ══════════════════════════════════════════════════════════════════
# SLIDE 11: Conformance Deviations
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "Specification Deviations Found", [
    "Deviation 1: IPv4 headers with IHL encoding < 20 bytes accepted",
    "   → smoltcp allows sub-minimum headers if buffer is long enough",
    "   → Could allow crafted packets to bypass validation",
    "",
    "Deviation 2: UDP over IPv6 accepts checksum = 0",
    "   → RFC 2460 mandates non-zero UDP checksums on IPv6",
    "   → Could mask corrupted datagrams on IPv6 networks",
    "",
    "Both are deliberate design choices documented in the codebase",
])


# ══════════════════════════════════════════════════════════════════
# SLIDE 12: Performance
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_image_slide(slide, "Performance: Loopback Throughput",
                os.path.join(FIGURES, 'fig4_performance_throughput.png'),
                "WSL avg: 52.15 Gbps | Windows avg: 17.11 Gbps | 3× performance gap")


# ══════════════════════════════════════════════════════════════════
# SLIDE 13: Microbenchmarks
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_image_slide(slide, "Microbenchmark Results",
                os.path.join(FIGURES, 'fig10_microbench.png'),
                "ipv4_parse: ~11ns | udp_parse_emit: ~24ns | ring_buffer_cycle: <1ns")


# ══════════════════════════════════════════════════════════════════
# SLIDE 14: White-Box Coverage
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_image_slide(slide, "White-Box Coverage Analysis",
                os.path.join(FIGURES, 'fig5_whitebox_coverage.png'),
                "Overall: 80.63% regions | assembler.rs: 98.28% | tcp.rs needs improvement: 72.58%")


# ══════════════════════════════════════════════════════════════════
# SLIDE 15: Enhanced Coverage
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_image_slide(slide, "Enhanced Suite Coverage",
                os.path.join(FIGURES, 'fig11_enhanced_coverage.png'),
                "IPv4 Combinatorial: 98.53% | Microbench: 98.76% | Only 6 total missed regions")


# ══════════════════════════════════════════════════════════════════
# SLIDE 16: Safety
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "Software Safety Assessment", [
    "Panic-safety tests: all passed on both platforms",
    "18 unsafe code occurrences in 4 files",
    "All unsafe blocks confined to FFI/platform-specific PHY layer",
    "Zero unsafe code in protocol parsing or state machine logic",
    "tuntap_interface.rs: 6 | bpf.rs: 5 | raw_socket.rs: 5 | mod.rs: 2",
], os.path.join(FIGURES, 'fig8_unsafe_inventory.png'))


# ══════════════════════════════════════════════════════════════════
# SLIDE 17: Portability
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_image_slide(slide, "Cross-Platform Portability",
                os.path.join(FIGURES, 'fig14_portability_matrix.png'),
                "8/9 suites passed on both platforms | Fuzzing blocked on Windows (MSVC/ASAN)")


# ══════════════════════════════════════════════════════════════════
# SLIDE 18: Quality Dashboard
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_image_slide(slide, "Quality Assessment Dashboard",
                os.path.join(FIGURES, 'fig12_test_dashboard.png'),
                "Mutation adequacy (69.67%) is the primary area for improvement")


# ══════════════════════════════════════════════════════════════════
# SLIDE 19: Key Findings
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "Key Findings", [
    "✓ 333 mutants tested — 69.18% adjusted kill ratio",
    "✓ 696 functional test cases — 100% pass rate",
    "✓ Zero crashes from coverage-guided fuzzing (Rust safety)",
    "✓ 2 RFC deviations documented (IPv4 header len, IPv6 UDP checksum)",
    "✓ 3× performance gap between WSL and Windows",
    "✓ 98.5%+ coverage for enhanced test suites",
    "✓ 18 unsafe occurrences — all in FFI layer, none in protocol logic",
    "✗ Windows fuzzing blocked by MSVC/ASAN toolchain issue",
])


# ══════════════════════════════════════════════════════════════════
# SLIDE 20: Recommendations
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "Recommendations", [
    "1. Add targeted tests for 102 surviving non-equivalent mutants",
    "2. Resolve Windows fuzzing via LLVM/Clang build configuration",
    "3. Enforce strict IPv6 UDP checksum validation via feature flag",
    "4. Extend fuzzing campaigns to ≥3600s per target",
    "5. Improve tcp.rs coverage (currently 72.58%) with state transition tests",
    "6. Consider address sanitizer testing on Linux for FFI code",
    "7. Automate cross-platform CI pipelines for both Windows and Linux",
])


# ══════════════════════════════════════════════════════════════════
# SLIDE 21: Thank You
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_BLUE
shape.line.fill.background()

add_text(slide, "Thank You", 1, 2, 11, 1,
         font_size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "Questions & Discussion", 1, 3.2, 11, 0.8,
         font_size=28, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, "All test artifacts, logs, and source code are available in the repository", 1, 4.5, 11, 0.6,
         font_size=18, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, "security_testing/ | security_testing/additions/ | security_testing/reports/", 1, 5.2, 11, 0.6,
         font_size=16, color=GRAY, align=PP_ALIGN.CENTER)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.35), prs.slide_width, Inches(0.15))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_BLUE
shape.line.fill.background()


# ── Save ──
prs.save(OUT)
print(f"Presentation saved to: {OUT}")
print(f"Total slides: {len(prs.slides)}")
